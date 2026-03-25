"""Text extraction from various file types for content classification.

All library imports are optional - extraction gracefully returns empty string
if the required package is not installed.
"""

from __future__ import annotations

import logging
import re
from pathlib import Path

from src.core.config import MAX_EXTRACT_SIZE
from src.core.models import FileInfo, FileType

logger = logging.getLogger(__name__)

# Max chars to extract for classification (keep prompts small)
MAX_TEXT_CHARS = 3000

# ── Check which optional packages are available ──
_HAS_PYMUPDF = False
try:
    import pymupdf  # noqa: F401
    _HAS_PYMUPDF = True
except ImportError:
    try:
        import fitz  # noqa: F401 - older pymupdf import name
        _HAS_PYMUPDF = True
    except ImportError:
        pass

_HAS_DOCX = False
try:
    import docx  # noqa: F401
    _HAS_DOCX = True
except ImportError:
    pass

_HAS_OPENPYXL = False
try:
    import openpyxl  # noqa: F401
    _HAS_OPENPYXL = True
except ImportError:
    pass

_HAS_PPTX = False
try:
    import pptx  # noqa: F401
    _HAS_PPTX = True
except ImportError:
    pass

_HAS_PILLOW = False
try:
    from PIL import Image  # noqa: F401
    _HAS_PILLOW = True
except ImportError:
    pass

_HAS_CHARDET = False
try:
    import chardet  # noqa: F401
    _HAS_CHARDET = True
except ImportError:
    pass


def get_available_extractors() -> dict[str, bool]:
    """Return which extractors are available (for UI status display)."""
    return {
        "PDF (PyMuPDF)": _HAS_PYMUPDF,
        "Word (python-docx)": _HAS_DOCX,
        "Excel (openpyxl)": _HAS_OPENPYXL,
        "PowerPoint (python-pptx)": _HAS_PPTX,
        "Images/EXIF (Pillow)": _HAS_PILLOW,
        "Encoding detection (chardet)": _HAS_CHARDET,
    }


def extract_text(file_info: FileInfo, max_size: int = MAX_EXTRACT_SIZE) -> str:
    """Extract text content from a file for classification. Returns up to MAX_TEXT_CHARS."""
    if file_info.size_bytes > max_size:
        return ""

    extractors: dict[FileType, tuple[callable, bool]] = {
        FileType.PDF: (_extract_pdf, _HAS_PYMUPDF),
        FileType.WORD: (_extract_docx, _HAS_DOCX),
        FileType.EXCEL: (_extract_excel, True),  # CSV fallback always works
        FileType.POWERPOINT: (_extract_pptx, _HAS_PPTX),
        FileType.TEXT: (_extract_text_file, True),  # Always available
    }

    entry = extractors.get(file_info.file_type)
    if not entry:
        return ""

    extractor_fn, is_available = entry
    if not is_available:
        logger.info(
            f"Skipping text extraction for {file_info.name} "
            f"(missing library for {file_info.file_type.value})"
        )
        return ""

    try:
        text = extractor_fn(file_info.path)
        # Clean and truncate
        text = _clean_text(text)
        file_info.content_preview = text[:500]
        return text[:MAX_TEXT_CHARS]
    except Exception as e:
        logger.warning(f"Failed to extract text from {file_info.path}: {e}")
        return ""


def extract_image_metadata(file_info: FileInfo) -> dict:
    """Extract EXIF and other metadata from images."""
    if not _HAS_PILLOW:
        return {}

    metadata = {}
    try:
        from PIL import Image
        from PIL.ExifTags import TAGS

        with Image.open(file_info.path) as img:
            metadata["width"] = img.width
            metadata["height"] = img.height
            metadata["format"] = img.format

            exif_data = img.getexif()
            if exif_data:
                for tag_id, value in exif_data.items():
                    tag_name = TAGS.get(tag_id, tag_id)
                    if isinstance(value, (str, int, float)):
                        metadata[str(tag_name)] = value
    except Exception as e:
        logger.debug(f"No EXIF data for {file_info.path}: {e}")

    return metadata


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using PyMuPDF."""
    # pymupdf >= 1.24 uses `import pymupdf`, older versions use `import fitz`
    try:
        import pymupdf as pdf_lib
    except ImportError:
        import fitz as pdf_lib

    text_parts = []
    with pdf_lib.open(str(path)) as doc:
        # Read first 10 pages max for classification
        for page_num in range(min(len(doc), 10)):
            page = doc[page_num]
            text_parts.append(page.get_text())

    return "\n".join(text_parts)


def _extract_docx(path: Path) -> str:
    """Extract text from Word documents."""
    from docx import Document

    doc = Document(str(path))
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def _extract_excel(path: Path) -> str:
    """Extract text from Excel files (sheet names + first rows)."""
    ext = path.suffix.lower()

    if ext in (".csv", ".tsv"):
        return _extract_text_file(path)

    if not _HAS_OPENPYXL:
        logger.info(f"Skipping Excel extraction for {path.name} (openpyxl not installed)")
        return ""

    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    text_parts = []

    for sheet_name in wb.sheetnames[:5]:
        text_parts.append(f"Sheet: {sheet_name}")
        ws = wb[sheet_name]
        for row in ws.iter_rows(max_row=20, values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                text_parts.append(" | ".join(cells))

    wb.close()
    return "\n".join(text_parts)


def _extract_pptx(path: Path) -> str:
    """Extract text from PowerPoint files."""
    from pptx import Presentation

    prs = Presentation(str(path))
    text_parts = []

    for slide_num, slide in enumerate(prs.slides[:15], 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_text.append(para.text.strip())
        if slide_text:
            text_parts.append(f"Slide {slide_num}: {' '.join(slide_text)}")

    return "\n".join(text_parts)


def _extract_text_file(path: Path) -> str:
    """Extract text from plain text files with encoding detection."""
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        if _HAS_CHARDET:
            try:
                import chardet
                raw = path.read_bytes()[:10000]
                detected = chardet.detect(raw)
                encoding = detected.get("encoding", "latin-1") or "latin-1"
                return raw.decode(encoding, errors="replace")
            except Exception:
                pass
        # Final fallback - latin-1 never throws
        return path.read_bytes()[:10000].decode("latin-1", errors="replace")


def _clean_text(text: str) -> str:
    """Clean extracted text for classification."""
    # Remove excessive whitespace
    text = re.sub(r"\s+", " ", text)
    # Remove non-printable characters (keep Hindi/Devanagari)
    text = re.sub(r"[^\w\s@.,;:!?₹$%&*()\-/\\'\"\u0900-\u097F]", "", text)
    return text.strip()
